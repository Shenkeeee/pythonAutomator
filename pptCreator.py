import openpyxl
from pptx import Presentation
from pptx.util import Inches

# getting excel page
wb = openpyxl.load_workbook('D:/Dokumentumok/Ibolya/mátrix.xlsx')
ws = wb.active

# Determine the size of the table
num_rows = ws.max_row
num_cols = ws.max_column



# region Create a new PowerPoint presentation and slides
prs = Presentation()
prs.slide_width = Inches(16)
prs.slide_height = Inches(9)

# title
slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(slide_layout)

title1 = slide.shapes.title
subtitle = slide.placeholders[1]

title1.text = "Tervteljesítés"
subtitle.text = "sdfdsffxd"

slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)

# endregion

# Determine the position and size of the table on the slide
left = Inches(0.5)
top = Inches(1.5)

# using fixed values for now
width = Inches(15)
height = Inches(2)

# Create a new table on the slide
table = slide.shapes.add_table(num_rows, num_cols, left, top, width, height).table


# Copy the data from the Excel file into the table
for i, row in enumerate(ws.iter_rows(values_only=True)):
    for j, cell_value in enumerate(row):
        table.cell(i, j).text = str(cell_value)
# xd
# save
prs.save('D:\Dokumentumok\Ibolya\mátrix.pptx')
